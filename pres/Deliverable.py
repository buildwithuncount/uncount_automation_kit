from pptx import Presentation
from pptx.util import Cm
import pandas as pd
import os
import shutil

import config as config


class Deliverable:

    amount = 0

    ipnut_file = config.input_file
    input_folder = config.input_folder

    output_file = config.output_file
    output_folder = config.output_folder

    def __init__(self, output_file, output_folder, input_file, input_folder, df):

        shutil.copy2(f"{input_folder}/{input_file}", f"{output_folder}/{output_file}")

        self.presentation = Presentation(f"{output_folder}/{output_file}")

        self.name = config.output_folder
        self.age = config.output_file
        self.df = df

        Deliverable.amount+=1

    def fill_slides(self): 
        print("...")
        
class OneSecondDeck(Deliverable):

    output_file = config.output_file
    output_folder = config.output_folder

    def __init__(self, output_file, output_folder, input_file, input_folder, answers, df):
        super(OneSecondDeck, self).__init__(output_file, output_folder, input_file, input_folder, df)
        self.salary = 233

        self.answers = answers
    
    def createPres(self):

        self.presentation.save(f"{config.output_folder}/{config.output_file}")
        print(f"# # # # Datei erstellt # # # #")

        return  

    def insertTextOnSlide(self, content, slide_nr, platzhalter_name):

        slide_nr -= 1

        slide = self.presentation.slides[slide_nr]

        gesuchter_platzhalter = None

        for platzhalter in slide.placeholders:
            if platzhalter.name == platzhalter_name:
                gesuchter_platzhalter = platzhalter
                break

        if gesuchter_platzhalter:
            gesuchter_platzhalter.text = content
            print(f"Erfolgreich platziert: {platzhalter_name}")

        else:
            print(f"Platzhalter mit dem Namen '{platzhalter_name}' wurde nicht gefunden.")

        return True

    def insertImageOnSlide(self, image_path, slide_nr, placeholder_name, pos_left:int, pos_top:int, pos_width:int, pos_height:int, placeBehind = False):
        try:
            slide_nr -= 1  

            slide = self.presentation.slides[slide_nr]

            gesuchter_platzhalter = None
            for platzhalter in slide.placeholders:
                if platzhalter.name == placeholder_name:
                    gesuchter_platzhalter = platzhalter
                    break

            if gesuchter_platzhalter:
                # Bild in die Slide einfügen und die Größe in Zentimetern anpassen
                left = Cm(pos_left)  # Horizontaler Abstand vom linken Rand (1 Zoll entspricht 2.54 Zentimetern)
                top = Cm(pos_top)   # Vertikaler Abstand vom oberen Rand (1 Zoll entspricht 2.54 Zentimetern)
                width = Cm(pos_width)  # Breite des Bilds (3 Zoll entsprechen 7.62 Zentimetern)
                height = Cm(pos_height) # Höhe des Bilds (2 Zoll entsprechen 5.08 Zentimetern)
                picture = slide.shapes.add_picture(image_path, left, top, width, height)
                
                if placeBehind == True:
                    shapes = slide.shapes._spTree
                    shapes.remove(picture._element)
                    shapes.insert(2, picture._element)  # Füge das Bild an Position 2 ein, um es unter anderen Objekten zu platzieren
                    
                print(f"Erfolgreich Bild in {placeholder_name} auf Slide {slide_nr + 1} in Zentimetern eingefügt.")
            else:
                print(f"Platzhalter mit dem Namen '{placeholder_name}' wurde nicht gefunden auf Slide {slide_nr + 1}.")
        except Exception as e:
            print(f"Fehler beim Einfügen des Bilds: {str(e)}")

    def insertChartOnSlide(self, df, func, slide, placeholder, pos_left, pos_top, pos_width, pos_height):
        chart_image, filtered_df = func(df) 
        self.insertImageOnSlide(chart_image, slide, placeholder, pos_left, pos_top, pos_width, pos_height)
        chart_image.close()
        return filtered_df

    def writeNotes(self, content, slide_nr):

        slide_nr -= 1

        slide = self.presentation.slides[slide_nr]

        notes_slide = slide.notes_slide
        
        notes_slide.notes_text_frame.text = content
        return True

print(config.title)
