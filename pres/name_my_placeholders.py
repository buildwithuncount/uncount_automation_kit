from pptx import Presentation

# Name ohne.pptx
file_name = "magic_report"

# Pfad zur Präsentation
pptx_path = f"./masters/{file_name}/{file_name}.pptx"

# Lade die Präsentation
presentation = Presentation(pptx_path)

# Schleife durch alle Slides in der Präsentation
for slide in presentation.slides:
    # Schleife durch alle Platzhalter auf der aktuellen Slide
    for platzhalter in slide.placeholders:
        # Platzhalter-Namen ermitteln
        platzhalter_name = platzhalter.name
        # Den Platzhalter mit seinem eigenen Namen beschriften
        platzhalter.text = platzhalter_name
        
for slide in presentation.slides:
    master_slide = slide.slide_layout

    # Informationen über den Master-Slide abrufen
    master_title = master_slide.name
    notes_slide = slide.notes_slide
    text_box = notes_slide.notes_text_frame
    text_box.text = f"Verwendeter Master-Slide: {master_title}"
    #print(master_title)

# Die geänderte Präsentation speichern
presentation.save(f'./masters/{file_name}/{file_name}_with_Names.pptx')